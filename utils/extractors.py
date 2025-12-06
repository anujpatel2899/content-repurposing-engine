"""Content extraction utilities (URL and File)."""
import requests
import io
import tempfile
import os
from typing import Tuple, Dict
import google.generativeai as genai


def extract_from_url(url: str) -> Tuple[str, str]:
    """
    Extract content from URL using Jina AI Reader (primary) or BeautifulSoup (fallback).
    
    Returns:
        (extracted_text, method_used)
    """
    # Try Jina AI Reader first
    try:
        jina_url = f"https://r.jina.ai/{url}"
        response = requests.get(jina_url, timeout=15)
        response.raise_for_status()
        
        text = response.text
        if text and len(text) > 100:
            return text, "Jina AI Reader"
    except Exception as e:
        print(f"Jina AI failed: {e}, trying BeautifulSoup...")
    
    # Fallback to BeautifulSoup
    try:
        from bs4 import BeautifulSoup
        
        response = requests.get(url, timeout=10)
        response.raise_for_status()
        soup = BeautifulSoup(response.content, 'html.parser')
        
        # Remove unwanted elements
        for script in soup(["script", "style", "nav", "footer", "header"]):
            script.decompose()
        
        text = soup.get_text(separator='\n')
        
        # Clean whitespace
        lines = (line.strip() for line in text.splitlines())
        chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
        text = '\n'.join(chunk for chunk in chunks if chunk)
        
        return text, "BeautifulSoup"
    except Exception as e:
        raise Exception(f"URL extraction failed: {str(e)}")


def extract_from_file(file_bytes: bytes, filename: str, google_api_key: str = None) -> Tuple[str, str]:
    """
    Extract content from uploaded file.
    
    Supports: PDF, DOCX, PPTX, TXT, MD
    
    Returns:
        (extracted_text, method_used)
    """
    file_ext = filename.split('.')[-1].lower()
    
    # PDF, DOCX, PPTX - Try Gemini first if key provided
    if file_ext in ['pdf', 'docx', 'pptx'] and google_api_key:
        try:
            genai.configure(api_key=google_api_key)
            
            # Create temp file
            with tempfile.NamedTemporaryFile(delete=False, suffix=f".{file_ext}") as tmp:
                tmp.write(file_bytes)
                tmp_path = tmp.name
            
            try:
                # Upload to Gemini
                uploaded_file = genai.upload_file(path=tmp_path)
                
                # Wait for processing
                import time
                while uploaded_file.state.name == "PROCESSING":
                    time.sleep(1)
                    uploaded_file = genai.get_file(uploaded_file.name)
                
                if uploaded_file.state.name == "FAILED":
                    raise Exception("Gemini file processing failed")
                
                # Extract text
                model = genai.GenerativeModel('gemini-2.0-flash-lite')
                response = model.generate_content([
                    "Extract all text from this document. Preserve structure. Return only text.",
                    uploaded_file
                ])
                
                # Cleanup
                os.unlink(tmp_path)
                genai.delete_file(uploaded_file.name)
                
                return response.text, "Gemini File API"
            except Exception as e:
                if os.path.exists(tmp_path):
                    os.unlink(tmp_path)
                # Fall through to direct extraction
                print(f"Gemini failed: {e}, using fallback...")
        except Exception as e:
            print(f"Gemini processing error: {e}")
    
    # Fallback: Direct extraction
    try:
        if file_ext == 'pdf':
            import pypdf
            pdf_file = io.BytesIO(file_bytes)
            reader = pypdf.PdfReader(pdf_file)
            text = ""
            for page in reader.pages:
                text += page.extract_text() + "\n"
            return text, "PyPDF"
        
        elif file_ext == 'docx':
            import docx
            doc = docx.Document(io.BytesIO(file_bytes))
            text = '\n'.join([para.text for para in doc.paragraphs])
            return text, "python-docx"
        
        elif file_ext == 'pptx':
            from pptx import Presentation
            prs = Presentation(io.BytesIO(file_bytes))
            full_text = []
            
            for i, slide in enumerate(prs.slides):
                full_text.append(f"\n--- Slide {i+1} ---\n")
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text = shape.text.strip()
                        if text:
                            full_text.append(text)
            
            return '\n'.join(full_text), "python-pptx"
        
        elif file_ext in ['txt', 'md']:
            text = file_bytes.decode('utf-8')
            return text, "Direct Read"
        
        else:
            raise Exception(f"Unsupported file type: {file_ext}")
    
    except Exception as e:
        raise Exception(f"File extraction failed: {str(e)}")
